"""Loader документов реального корпуса «Научный клубок».

Поддерживает форматы: PDF · DOCX · DOC · PPTX · PPT · XLSX · XLS · TXT.
Старые бинарные .doc/.ppt/.xls конвертируются через libreoffice, если доступен.

Обходит дерево источников рекурсивно и извлекает метаданные из пути:
  Источники информации/
    ├── Доклады/…                       → doc_type = report
    ├── Журналы/<Издание>/<Год>/…       → doc_type = journal, journal, year
    ├── Статьи/…                        → doc_type = article
    ├── Обзоры/…                        → doc_type = review
    └── Материалы конференций/…         → doc_type = conference

Также определяет язык (RU/EN по кириллице) и country_code — для гео-фильтра.
"""

from __future__ import annotations

import hashlib
import re
import subprocess
import tempfile
import uuid
from pathlib import Path

from loguru import logger

from app.config import settings
from app.services import geo

_TYPE_KEYWORDS = [
    ("материалы конференц", "conference"),
    ("конференц", "conference"),
    ("доклад", "report"),
    ("презентац", "report"),
    ("журнал", "journal"),
    ("обзор", "review"),
    ("стать", "article"),
    ("патент", "patent"),
    ("диссертац", "thesis"),
]

_YEAR_RE = re.compile(r"(?<!\d)(19|20)\d{2}(?!\d)")


def classify_path(path, root=None) -> dict:
    """Извлекает doc_type / source_category / journal / year из пути документа."""
    p = Path(path)
    if root:
        try:
            rel = p.relative_to(root)
        except ValueError:
            rel = p
    else:
        rel = p
    parts = list(rel.parts)
    low = [seg.lower() for seg in parts]

    doc_type = "document"
    category = None
    for i, seg in enumerate(low):
        for kw, dtype in _TYPE_KEYWORDS:
            if kw in seg:
                doc_type = dtype
                category = parts[i]
                break
        if category:
            break

    journal = None
    for i, seg in enumerate(low):
        if "журнал" in seg and i + 1 < len(parts):
            journal = parts[i + 1]
            break

    year = None
    for seg in parts:
        m = _YEAR_RE.search(seg)
        if m:
            year = int(m.group(0))
            break

    return {
        "doc_type": doc_type,
        "source_category": category,
        "journal": journal,
        "year": year,
    }


def file_hash(path):
    h = hashlib.sha1()
    with open(path, "rb") as f:
        for chunk in iter(lambda: f.read(65536), b""):
            h.update(chunk)
    return h.hexdigest()


def stable_doc_id(seed):
    return "DOC-" + hashlib.sha1(str(seed).encode("utf-8")).hexdigest()[:12].upper()


def detect_language(text):
    if not text:
        return None
    sample = text[:1000].lower()
    cyr = sum(1 for c in sample if "а" <= c <= "я")
    lat = sum(1 for c in sample if "a" <= c <= "z")
    if cyr > lat * 0.7:
        return "ru"
    if lat > 0:
        return "en"
    return None


def _which(name):
    from shutil import which
    return which(name)


def extract_text_pdf(path):
    try:
        import pdfplumber
    except ImportError:
        logger.warning("pdfplumber not installed")
        return []
    pages = []
    try:
        with pdfplumber.open(path) as pdf:
            for i, page in enumerate(pdf.pages, 1):
                text = page.extract_text() or ""
                pages.append((i, text))
    except Exception as e:  # noqa: BLE001
        logger.warning(f"PDF parse failed for {path}: {e}")
    return pages


def extract_text_docx(path):
    try:
        from docx import Document
    except ImportError:
        return []
    try:
        d = Document(str(path))
        full = "\n".join(p.text for p in d.paragraphs if p.text.strip())
        for tbl in d.tables:
            for row in tbl.rows:
                cells = [c.text.strip() for c in row.cells if c.text.strip()]
                if cells:
                    full += "\n" + " | ".join(cells)
        return [(1, full)] if full else []
    except Exception as e:  # noqa: BLE001
        logger.warning(f"DOCX parse failed for {path}: {e}")
        return []


def extract_text_pptx(path):
    """Текст со слайдов: фигуры, таблицы и заметки докладчика."""
    try:
        from pptx import Presentation
    except ImportError:
        logger.warning("python-pptx not installed")
        return []
    try:
        prs = Presentation(str(path))
    except Exception as e:  # noqa: BLE001
        logger.warning(f"PPTX parse failed for {path}: {e}")
        return []
    pages = []
    for idx, slide in enumerate(prs.slides, 1):
        buf = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                buf.append(shape.text_frame.text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells if c.text.strip()]
                    if cells:
                        buf.append(" | ".join(cells))
        try:
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                note = slide.notes_slide.notes_text_frame.text.strip()
                if note:
                    buf.append("[заметки] " + note)
        except Exception:  # noqa: BLE001
            pass
        text = "\n".join(buf).strip()
        if text:
            pages.append((idx, text))
    return pages


def extract_text_excel(path):
    """Каждый лист → отдельная «страница»; строки склеиваются в текст."""
    try:
        import pandas as pd
    except ImportError:
        return []
    try:
        xl = pd.ExcelFile(path)
    except Exception as e:  # noqa: BLE001
        logger.warning(f"Excel open failed for {path}: {e}")
        return []
    pages = []
    for idx, sheet in enumerate(xl.sheet_names, 1):
        try:
            df = xl.parse(sheet, header=None, dtype=str)
        except Exception as e:  # noqa: BLE001
            logger.warning(f"Excel sheet '{sheet}' failed in {path}: {e}")
            continue
        lines = [f"[лист] {sheet}"]
        for _, row in df.iterrows():
            cells = [str(c).strip() for c in row.tolist()
                     if str(c).strip() and str(c).lower() != "nan"]
            if cells:
                lines.append(" | ".join(cells))
        text = "\n".join(lines).strip()
        if len(text) > len(sheet) + 8:
            pages.append((idx, text))
    return pages


def _soffice_convert(path: Path, target: str):
    if not settings.soffice_convert or not _which(settings.soffice_bin):
        return None
    outdir = Path(tempfile.mkdtemp(prefix="soffice_"))
    try:
        subprocess.run(
            [settings.soffice_bin, "--headless", "--convert-to", target,
             "--outdir", str(outdir), str(path)],
            check=True, capture_output=True, timeout=180,
        )
    except Exception as e:  # noqa: BLE001
        logger.warning(f"libreoffice convert failed for {path.name}: {e}")
        return None
    ext = target.split(":")[0]
    out = outdir / (path.stem + "." + ext)
    return out if out.exists() else None


def load_document(path, root=None):
    p = Path(path)
    suffix = p.suffix.lower()

    if suffix == ".pdf":
        pages = extract_text_pdf(p)
    elif suffix in (".docx", ".docm"):
        pages = extract_text_docx(p)
    elif suffix == ".doc":
        conv = _soffice_convert(p, "docx")
        pages = extract_text_docx(conv) if conv else []
    elif suffix == ".pptx":
        pages = extract_text_pptx(p)
    elif suffix == ".ppt":
        conv = _soffice_convert(p, "pptx")
        pages = extract_text_pptx(conv) if conv else []
    elif suffix in (".xlsx", ".xls"):
        pages = extract_text_excel(p)
    elif suffix == ".txt":
        pages = [(1, p.read_text(encoding="utf-8", errors="ignore"))]
    else:
        logger.warning(f"Unsupported document type: {p}")
        pages = []

    full_text = " ".join(t for _, t in pages)
    lang = detect_language(full_text)
    country = geo.guess_country_from_text(full_text, language=lang)
    region = geo.classify_region(country)
    meta = classify_path(p, root)

    return {
        "doc_id": stable_doc_id(file_hash(p)),
        "file_path": str(p),
        "title": p.stem,
        "kind": "file",
        "format": suffix.lstrip("."),
        "language": lang,
        "country_code": country,
        "geo_region": region,
        "doc_type": meta["doc_type"],
        "source_category": meta["source_category"],
        "journal": meta["journal"],
        "year": meta["year"],
        "page_count": len(pages),
        "pages": pages,
    }


def chunk_text(text, size=None, overlap=None):
    size = size or settings.chunk_size
    overlap = overlap or settings.chunk_overlap
    text = re.sub(r"\s+", " ", text or "").strip()
    if not text:
        return []
    chunks, start = [], 0
    while start < len(text):
        end = min(start + size, len(text))
        chunks.append(text[start:end])
        if end == len(text):
            break
        start += size - overlap
    return chunks


def iter_chunks(document):
    for page_num, page_text in document["pages"]:
        for chunk in chunk_text(page_text):
            yield {
                "chunk_id": str(uuid.uuid4()),
                "doc_id": document["doc_id"],
                "page": page_num,
                "text": chunk,
            }


def _excluded(path: Path, exclude_dirs) -> bool:
    parts_low = [seg.lower() for seg in path.parts]
    return any(d.lower() in parts_low for d in exclude_dirs)


def iter_document_paths(root, exclude_dirs=None):
    root = Path(root)
    exclude_dirs = exclude_dirs if exclude_dirs is not None else settings.corpus_exclude_dirs
    exts = set(settings.doc_extensions)
    if not root.exists():
        return
    for f in sorted(root.rglob("*")):
        if not f.is_file() or f.suffix.lower() not in exts:
            continue
        if _excluded(f, exclude_dirs):
            continue
        yield f


def iter_documents(corpus_dir, root=None):
    corpus = Path(corpus_dir)
    walk_root = Path(root) if root else corpus

    legacy = corpus / "documents"
    if legacy.exists() and not settings.corpus_recursive:
        walk_root = legacy

    for f in iter_document_paths(walk_root):
        try:
            yield load_document(f, root=walk_root)
        except Exception as e:  # noqa: BLE001
            logger.warning(f"failed to load {f}: {e}")


# --- Top-level worker: должен быть picklable, поэтому вынесен из класса. ---
def _parse_worker(args):
    path_str, root_str = args
    try:
        return load_document(Path(path_str), root=Path(root_str) if root_str else None)
    except Exception as e:  # noqa: BLE001
        return {"__error__": f"{type(e).__name__}: {e}", "__path__": path_str}


def iter_documents_parallel(corpus_dir, root=None, max_workers=None, skip_paths=None):
    """Параллельный обход: N процессов парсят PDF/DOCX/…, main получает
    словари документов через as_completed по мере готовности. Порядок
    не гарантируется, но нам всё равно — эмбеддинг и upsert идемпотентны.

    skip_paths — множество file_path, которые не нужно переспарсить
    (например, уже есть в Neo4j). Даёт resume после краша.
    """
    from concurrent.futures import ProcessPoolExecutor, as_completed
    import os as _os

    corpus = Path(corpus_dir)
    walk_root = Path(root) if root else corpus
    legacy = corpus / "documents"
    if legacy.exists() and not settings.corpus_recursive:
        walk_root = legacy

    all_paths = [str(f) for f in iter_document_paths(walk_root)]
    if not all_paths:
        return
    if skip_paths:
        before = len(all_paths)
        paths = [p for p in all_paths if p not in skip_paths]
        logger.info(f"Skip {before - len(paths)} already-ingested docs "
                    f"({len(paths)} to process)")
    else:
        paths = all_paths
    if not paths:
        return

    if max_workers is None:
        cpu = _os.cpu_count() or 4
        max_workers = max(2, min(8, cpu - 1))
    root_str = str(walk_root)

    # Крупные PDF (>40 MB) грузим меньшим числом воркеров — pdfplumber
    # держит ~1.5 GB на 100 MB PDF, 6 параллельных → OOM в контейнере.
    LARGE_BYTES = 40 * 1024 * 1024
    small_paths, large_paths = [], []
    for p in paths:
        try:
            (large_paths if _os.path.getsize(p) > LARGE_BYTES else small_paths).append(p)
        except OSError:
            small_paths.append(p)
    large_workers = max(2, max_workers // 2)

    def _run_pool(batch, workers, label):
        if not batch:
            return
        logger.info(f"Parsing {len(batch)} {label} docs with {workers} workers")
        with ProcessPoolExecutor(max_workers=workers) as pool:
            futures = [pool.submit(_parse_worker, (p, root_str)) for p in batch]
            for fut in as_completed(futures):
                try:
                    doc = fut.result()
                except Exception as e:  # noqa: BLE001
                    logger.warning(f"worker crashed: {e}")
                    continue
                if not doc:
                    continue
                if "__error__" in doc:
                    logger.warning(f"failed to load {doc.get('__path__')}: {doc['__error__']}")
                    continue
                yield doc

    yield from _run_pool(small_paths, max_workers, "small")
    yield from _run_pool(large_paths, large_workers, "large")
